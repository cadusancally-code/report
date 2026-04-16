import matplotlib.pyplot as plt
from pptx import Presentation


def export_excel(df):
    file = "relatorio.xlsx"
    df.to_excel(file, index=False)
    return file


def export_png(df):
    file = "grafico.png"

    fig, ax = plt.subplots()
    df.groupby('data')[['recebidas', 'tratadas']].sum().plot(ax=ax)
    plt.savefig(file)

    return file


def export_ppt(df):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    total_r = df['recebidas'].sum()
    total_t = df['tratadas'].sum()

    slide.shapes.title.text = "Relatório Operacional"
    slide.placeholders[1].text = f"Recebidas: {total_r}\nTratadas: {total_t}"

    file = "relatorio.pptx"
    prs.save(file)

    return file